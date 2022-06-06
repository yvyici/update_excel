import openpyxl
import re
#path = "hyb3c.xlsx"
path = "WIP.xlsx"
path_new="hyb3c_new.xlsx"
wb = openpyxl.load_workbook(path)


#sht2 = wb["總表"]
sht2 = wb["110nm Silvo FE_5V"]

from pptx import Presentation
path_to_presentation='Reference_lot.pptx'
path_to_presentation_new='Reference_lot_new.pptx'
prs = Presentation(path_to_presentation)
for slide in prs.slides:
    # print(slide)
    for shape in slide.shapes:
        # print(shape.has_table)
        if not shape.has_table:
            continue
        table=shape.table
        #print('table.cell(0,0).text===',table.cell(0,0).text)
pptrow = []
pptcol = []
for row_idx, row in enumerate(table.rows):
    pptrow.append(row_idx)
    for col_idx, cell in enumerate(row.cells):
        if col_idx not in pptcol:
            pptcol.append(col_idx)


## get the title in ppt which column
def pptcolnum(title):
    j = 0
    while j<10:
        # sc = sht1.cell(row= 1, column= j).value
        sc = table.cell(0,j).text
        # if type(sc) == str and title in sc:
        if title in sc:
            break;
        else:
            j=j+1
    return j

pptcolnum_LotID=pptcolnum('Lot ID')
pptcolnum_Product=pptcolnum('Product')
pptcolnum_pcs=pptcolnum('pcs')
pptcolnum_Status=pptcolnum('Status')

shtcolnum_LotID=5
shtcolnum_Product=3
shtcolnum_pcs=6


shtrownum_Layer=3

## get all ppt new id
def pptid():
    a = []
    i = 0
    while i < len(pptrow):
        sc = table.cell(i, pptcolnum_LotID).text
        print('kk sc===',sc)
        if type(sc) == str and re.search(r'(\S{7})', sc):
            a.append(re.findall(r'(\S{7})', sc)[0])
        else:
            if i<1:
                a.append(0)
        i= i+1
    #print(a)
    #print(ZZZ)
    return a

def unmerge():
    i = 1
    a = []
    for i in range(0,len(pptrow)):
        sc = table.cell(i, pptcolnum_Product+1)
        if re.search(r'\S+',sc.text):
            a.append(i)
            # print(i, True)
        elif i == len(pptrow) and re.search(r'\S+',sc.text):
            a.append(i)
        else:
            # print(i, False)
            pass
        i = i+1
    j = 0
    for j in range(0, len(a)-1):
        if a[j] == a[j+1]-1:
            pass
        else:
            sc = table.cell(a[j], pptcolnum_Product+1)
            sc.split()
        j=j+1
    for c in range(1, len(pptrow)):
        sc = table.cell(c, pptcolnum_Product+1)
        if re.search(r'\S+',sc.text) == None:
            sc.text = table.cell(c-1, pptcolnum_Product+1).text
            # fonts.size = Pt(10)
        # elif sht1:
        #     pass
        else:
            pass
        c=c+1

## get same and different id
def sameID():
    j= 0
    same= []
    a = []
    b = []
    c = []
    d = []
    pptid2=pptid()
    while j < len(pptid2):
        vl = pptid2[j]
        i= 7
        #print('sht2.max_row===',sht2.max_row)
        while i<=sht2.max_row:
            sc = sht2.cell(column=shtcolnum_LotID, row= i).value
            #print('sc===',sc)
            if vl == sc:
                same.append(sc)
                a.append(j)
                b.append(i)
                break;
            i = i+1
        if j not in a:
            c.append(j)
            d.append(vl)
        j = j+1
    print('same===',same)
    print('a===',a)
    print('b===',b)
    print('c===',c)   
    print('d===',d)
    #print(ZZZ)
    return same, a, b, c, d
# print(sameID()[0])
# print(sameID()[1])
# print(sameID()[2])
# print(sameID()[3])           #[0]same id [1]ppt column [2]總表 column [3]ppt not the same column
# print(sameID()[4])



# the ppt table build to dictionary
def ppt2dict():
    pptdict = {}
    j = []
    for i in range(1, len(pptrow)):
        sc = table.cell(i, pptcolnum_LotID).text
        #print(sc)
        product = table.cell(i,pptcolnum_Product+1).text
        if sc=='':
            continue
        #print(product)
        pcs = table.cell(i, pptcolnum_pcs).text
        # print(pcs)
        status = table.cell(i, pptcolnum_Status).text
        #print('status===',status)
        v = sc.strip()
        s = re.findall(r'\w{2}', status)[0]
        p = re.sub(r'[()]','', product)
        if re.search(r'\S{3,}', p):
            p = re.findall(r'\S{3,}', p)[0]
        pptdict[v] = {'Product':p, 'pcs':pcs, 'Status':s}
        i = i+1
    print('pptdict===',pptdict)
    return pptdict

## to get the status site
def upcol(title):
    j = 9
    while j<50:
        sc = sht2.cell(row= shtrownum_Layer, column= j).value
        if type(sc) == str and title in sc:
            break;
        else:
            j=j+1
    return j

## update the same lot id's status
def Update_OldLotID():
    from datetime import datetime
    timeNow = datetime.now()
    i = 0
    sameID2=sameID()
    print('sameID2===',sameID2)
#    print(ZZZ)
    if len(sameID()[1]) == len(sameID2[2]):
        while i < len(sameID2[0]):
            # sc1 = sht1.cell(row= sameID()[1][i], column= pptcolnum('Status')).value
            sc1 = ppt2dict()[sameID2[0][i]]['Status']
            pat = re.findall(r'(\S{2})', sc1)
            # print(pat[0])
            if upcol(pat[0]) == 50:
                print('Can\'t find the layer name "%s"'%pat[0])
                pass
            elif pat and upcol(pat[0]) < 50:
                date = str(timeNow.month)+'/'+str(timeNow.day)
                print(date)
                sht2.cell(row= sameID2[2][i], column= upcol(pat[0])).value = str(timeNow.month)+'/'+str(timeNow.day)
            else:
                print('error')
            i = i+1
    else:
        print('Error the two elements is not match')


# Add New Product
def Add_NewProduct():
    i = 7
    pptid2=pptid()
    old = []
    while i<=sht2.max_row:
        sc = sht2.cell(column=shtcolnum_Product, row=i).value
        if type(sc) == str:
            p = re.sub(r'[()]','', sc)
            if re.search(r'\S{3,}', p):
                p = re.findall(r'\S{3,}', p)[0]
        if type(p) == str:
            if p not in old:
                old.append(p)
        i = i+1
    print('for old ', old)

    i = 0
    new = []
    for i in range(1, len(pptid2)):
        pr = table.cell(i, pptcolnum_Product+1).text
        p = re.sub(r'[()]','', pr)
        if re.search(r'\S{3,}', p):
            p = re.findall(r'\S{3,}', p)[0]
        if p not in new:
            new.append(p)
        else:
            pass
        i = i+1
    print('for new ', new)

    i = 0
    j = 0
    #print('old===',old)
    #print('new===',new)
    for i in range(0, len(new)):
        for j in range(0, len(old)):
            if new[i] in old[j]:
                print("It's already in.")
                break;
            elif j< len(old)-2 and new[i] not in old[j]:
                j = j+1
                pass
            elif j == len(old)-1 and new[i] not in old[j]:
                sht2.cell(row= sht2.max_row+2, column=2).value = new[i]
                print('Can\'t fint the 型號 %s, we\'ll add in the last row'%new[i])
                # sht2.cell(row= sht2.max_row, column=4).value = sht1.cell(row=3,column=pptcolnum('Lot ID')).value
                # sht2.cell(row= sht2.max_row, column=5).value = sht1.cell(row=3,column=pptcolnum('pcs')).value
        i = i+1


## tide_up_info
def getstart(kindname):
    i=7
    i0=[]
    i1=[]
    while i <= sht2.max_row:
        i1 = sht2.cell(column=shtcolnum_Product, row= i).value
        if type(i1) == str:
            if kindname in i1:
                i0.append(i-1)
            elif kindname not in i1 and len(i0)>0:
                break;
        i = i+1
    return(i-1)



## Add the new Lot ID
def Add_NewLotID(): # ne lot
    from datetime import datetime
    timeNow = datetime.now()
    i = 1
    j = 1
    sameID2=sameID()
    print('hello sameID()[4]===',sameID()[4])
    #print(ZZZ)
    mergedRanges = sht2.merged_cells.ranges
    i=0
    q = []
    while mergedRanges:
        for entry in mergedRanges:
            i+=1
            q.append(str(entry))
            print("  unMerging: " + str(i) + ": " +str(entry))
            sht2.unmerge_cells(str(entry))
    t = []
    while j < len(sameID2[4]):
        print(sameID2[4])
        print(j)
        p = sameID2[4]
        print('?'+str(p[j])+'?')
        pd = ppt2dict()[p[j]]
        print('pd===',pd)
        c = getstart(pd['Product'])
        print('hello c===',c)
        c=c
        t.append(c)
        #copy_style=False
        sht2.insert_rows(c,1)
        print('='+str(c)+'=')
        sht2.cell(row=c, column= shtcolnum_LotID).value = p[j]
        # sht2.cell(row=c, column= shtcolnum_Product).value = pd['Product']
        sht2.cell(row=c, column= shtcolnum_pcs).value = int(pd['pcs'])
        print('len(sameID()[4])===',len(sameID2[4]))
        print('sameID()[4]===',sameID2[4])
        for i in range(1, len(sameID2[4])):
            sc1 = pd['Status']
            pat = re.findall(r'(\S{2})', sc1)
            print('pat[0]===',pat[0])
            if upcol(pat[0]) == 50:
                print('Can\'t find the layer name "%s"'%pat[0])
                break;
            elif pat and upcol(pat[0]) < 50:
                date = str(timeNow.month)+'/'+str(timeNow.day)
                # print(date)
                sht2.cell(row= c, column= upcol(pat[0])).value = date
            else:
                print('error')
            i = i+1
        j=j+1
    print(t)

    l2 = []
    for k in range(7, sht2.max_row+1):
        c2 = sht2.cell(k, 2).value
        if type(c2) == str:
            l2.append(k)
        elif type(c2) != str:
            pass
        k=k+1

    l3 = []
    for k in range(7, sht2.max_row+1):
        c3 = sht2.cell(k, 3).value
        if type(c3) == str:
            l3.append(k)
        elif type(c3) != str:
            pass
        k=k+1

    print(l2,l3)

    for ll in range(0, len(l2)-1):
        sht2.merge_cells('B'+str(l2[ll])+':'+'B'+str(l2[ll+1]-1))
        # print('B'+str(l2[ll])+':'+'B'+str(l2[ll+1]-1))
        ll+=1

    for lll in range(0, len(l3)-1):
        sht2.merge_cells('C'+str(l3[lll])+':'+'C'+str(l3[lll+1]-1))
        # print('C'+str(l3[lll])+':'+'C'+str(l3[lll+1]-1))
        lll+=1

    sht2.merge_cells('B'+str(l2[len(l2)-1])+':'+'B'+str(sht2.max_row))
    # print('B'+str(l2[len(l2)-1])+':'+'B'+str(sht2.max_row))
    sht2.merge_cells('C'+str(l3[len(l3)-1])+':'+'C'+str(sht2.max_row))
    # print('C'+str(l3[len(l3)-1])+':'+'C'+str(sht2.max_row))
    
    # b = []
    # for g in range(len(q)):
    #     h = re.sub(r'\d+:\S+','', q[g])
    #     f = re.findall(r'\d+', q[g])
    #     b.append(h)
    #     # print(int(f[0]), int(f[1]))
    #     former = int(f[0])
    #     last0 = int(f[1])
    #     print(former, last0)
    #     if c <= last0 and g == 0:
    #         print(last0)
    #         last = last0+1
    #         print(last)
    #         # print(str(b[g])+str(former)+':'+str(b[g])+str(last))
    #     elif c > last0 and g == 0:
    #         print(last0)
    #         # print(str(b[g])+str(former)+':'+str(b[g])+str(last))
    #         pass
    #     elif c<=last0 and sht2.cell(last,3).value == 'Y1931A':
    #         print(last0)
    #         last= last0+1
    #         print(last)
    #         # print(str(b[g])+str(former)+':'+str(b[g])+str(last))
    #     elif c<=last0 and g>0:
    #         print(last0)
    #         former = former+1
    #         last = last0+1
    #         print(last)
    #         # print(str(b[g])+str(former)+':'+str(b[g])+str(last))
    #     else:
    #         print(last0)
    #         print('Error')
    #         # print(str(b[g])+str(former)+':'+str(b[g])+str(last))
    #     # sht2.merge_cells(str(b[g])+str(former)+':'+str(b[g])+str(last))
    #     g = g+1


if __name__ == '__main__':
    
    unmerge()
    Update_OldLotID()      # update old 
    Add_NewProduct()   # new Product (fine)
    Add_NewLotID()       # Add new Lot 

wb.save(path_new)
prs.save(path_to_presentation_new)
wb.close()
